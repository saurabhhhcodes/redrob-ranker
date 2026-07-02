#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def create_presentation():
    prs = Presentation()
    
    # Set slide dimensions to 16:9 widescreen (13.33 inches x 7.5 inches)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Design colors
    BG_COLOR = RGBColor(11, 15, 25) # Deep Space Blue
    ACCENT_PRIMARY = RGBColor(124, 58, 237) # Violet
    ACCENT_SECONDARY = RGBColor(16, 185, 129) # Emerald
    TEXT_PRIMARY = RGBColor(243, 244, 246) # Light grey
    TEXT_SECONDARY = RGBColor(156, 163, 175) # Muted grey
    
    # Helper to style a slide background
    def set_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    # Helper to add standard header
    def add_header(slide, title_text, category_text):
        # Category tag
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.73), Inches(0.4))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.name = 'Plus Jakarta Sans'
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = ACCENT_SECONDARY
        
        # Main Title
        tx_box2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.73), Inches(0.8))
        tf2 = tx_box2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.name = 'Outfit'
        p2.font.size = Pt(28)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_PRIMARY

    # Slide 1: Cover Slide
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    
    # Title & Subtitle in a single text box
    tx_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.0))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Redrob Candidate Discovery Engine"
    p.font.name = 'Outfit'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = ACCENT_SECONDARY
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "The Self-Healing Autonomous Development Workspace"
    p2.font.name = 'Outfit'
    p2.font.size = Pt(20)
    p2.font.color.rgb = TEXT_PRIMARY
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(15)
    
    p3 = tf.add_paragraph()
    p3.text = "Track 1: The AI Systems Architect: Reimagining Work\nSaurabh Kumar Bajpai • saurabhbajpai03@outlook.com"
    p3.font.name = 'Plus Jakarta Sans'
    p3.font.size = Pt(13)
    p3.font.color.rgb = TEXT_SECONDARY
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(40)

    # Slide 2: The Problem
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    add_header(slide, "The Problem", "Developer Friction at Scale")
    
    # Left column bullets
    tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.5))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    bullets = [
        "Developers spend less than 30% of their workday writing features.",
        "The majority of time is consumed by context-switching, searching documentation, and environment setup.",
        "Manual validation cycles (writing test suites, auditing UI, finding lints) create bottlenecks.",
        "Keyword-stuffed and fraudulent candidate resumes waste HR screening resources."
    ]
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = "✦  " + b
        p.font.name = 'Plus Jakarta Sans'
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(14)
        
    # Right column box
    tx_box_r = slide.shapes.add_textbox(Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.5))
    tf_r = tx_box_r.text_frame
    tf_r.word_wrap = True
    p_r = tf_r.paragraphs[0]
    p_r.text = "THE CORPORATE COST"
    p_r.font.name = 'Outfit'
    p_r.font.size = Pt(16)
    p_r.font.bold = True
    p_r.font.color.rgb = RGBColor(239, 68, 68) # Red
    
    p_r2 = tf_r.add_paragraph()
    p_r2.text = "\nFragile deployments, slow developer onboarding, and high ATS processing noise result in massive resource drains for scaling enterprises."
    p_r2.font.name = 'Plus Jakarta Sans'
    p_r2.font.size = Pt(15)
    p_r2.font.color.rgb = TEXT_SECONDARY

    # Slide 3: The Solution
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    add_header(slide, "The Solution", "Autonomous Workspace Integration")
    
    # Left column bullets
    tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.5))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    bullets = [
        "Reimagines developer setups into a fully containerized autonomous agent workspace.",
        "Features self-healing feedback loops: AI compiles code, runs test suites, lints, and corrects syntax errors autonomously.",
        "Integrates Chrome DevTools MCP for automated visual and accessibility auditing before commits.",
        "Employs deterministic logical filters to flag ATS anomalies and keyword stuffers automatically."
    ]
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = "✦  " + b
        p.font.name = 'Plus Jakarta Sans'
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(14)
        
    # Right column box
    tx_box_r = slide.shapes.add_textbox(Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.5))
    tf_r = tx_box_r.text_frame
    tf_r.word_wrap = True
    p_r = tf_r.paragraphs[0]
    p_r.text = "ARCHITECTURAL PARADIGM SHIFT"
    p_r.font.name = 'Outfit'
    p_r.font.size = Pt(16)
    p_r.font.bold = True
    p_r.font.color.rgb = ACCENT_SECONDARY
    
    p_r2 = tf_r.add_paragraph()
    p_r2.text = "\nShifts developer responsibilities from manual coding and debug loops to high-leverage architectural orchestration."
    p_r2.font.name = 'Plus Jakarta Sans'
    p_r2.font.size = Pt(15)
    p_r2.font.color.rgb = TEXT_SECONDARY

    # Slide 4: Dashboard - Top View
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    add_header(slide, "Recruiter Dashboard", "Candidates Discovery & Match View")
    
    # Image on left
    if os.path.exists("dashboard_top.png"):
        slide.shapes.add_picture("dashboard_top.png", Inches(0.8), Inches(1.8), width=Inches(6.0), height=Inches(4.5))
        
    # Text on right
    tx_box = slide.shapes.add_textbox(Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.5))
    tf = tx_box.text_frame
    tf.word_wrap = True
    bullets = [
        "Match Score Calibration: Real-time evaluation based on years of experience, current title, and technical skillset.",
        "Custom AI Reasoning: Provides custom, non-templated match summaries generated dynamically for each candidate.",
        "Interactive Candidate List: Displays rank, name, and alignment scores for rapid navigation."
    ]
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = "✦  " + b
        p.font.name = 'Plus Jakarta Sans'
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(20)

    # Slide 5: Dashboard - Bottom View
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    add_header(slide, "Telemetry & Skills", "Behavioral Signals & Verification")
    
    # Text on left
    tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.3), Inches(4.5))
    tf = tx_box.text_frame
    tf.word_wrap = True
    bullets = [
        "Redrob Behavioral Signals: Visualizes recruiter response rates, notice periods, and login activity.",
        "Availability Highlights: Prioritizes candidates with a buyout timeline of <= 30 days.",
        "Technical Competencies: Groups and colors skills based on explicit expert, advanced, or intermediate proficiency."
    ]
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = "✦  " + b
        p.font.name = 'Plus Jakarta Sans'
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(20)
        
    # Image on right
    if os.path.exists("dashboard_bottom.png"):
        slide.shapes.add_picture("dashboard_bottom.png", Inches(6.5), Inches(1.8), width=Inches(6.0), height=Inches(4.5))

    # Slide 6: System Architecture
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    add_header(slide, "System Architecture", "Self-Healing Developer Loop")
    
    # 4 flowchart boxes (orchestrator -> sandbox -> devtools -> git push)
    nodes = [
        {"title": "Orchestrator", "desc": "Agent coordinating actions"},
        {"title": "Docker Sandbox", "desc": "Runs code & test suites"},
        {"title": "Chrome DevTools", "desc": "Audits UI & accessibility"},
        {"title": "Git Push", "desc": "Pushes clean branches"}
    ]
    for idx, node in enumerate(nodes):
        left = Inches(0.8 + idx * 3.0)
        top = Inches(2.2)
        width = Inches(2.5)
        height = Inches(1.8)
        
        # Add box shape
        shape = slide.shapes.add_shape(1, left, top, width, height) # 1 = Rectangle
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(20, 27, 45)
        shape.line.color.rgb = ACCENT_PRIMARY
        shape.line.width = Pt(2)
        
        tf = shape.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = node["title"]
        p.font.name = 'Outfit'
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ACCENT_SECONDARY
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = "\n" + node["desc"]
        p2.font.name = 'Plus Jakarta Sans'
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_SECONDARY
        p2.alignment = PP_ALIGN.CENTER
        
        # Add arrows
        if idx < 3:
            arrow_box = slide.shapes.add_textbox(left + width, top + Inches(0.6), Inches(0.5), Inches(0.5))
            arrow_tf = arrow_box.text_frame
            arrow_p = arrow_tf.paragraphs[0]
            arrow_p.text = "➔"
            arrow_p.font.size = Pt(24)
            arrow_p.font.color.rgb = ACCENT_PRIMARY
            arrow_p.alignment = PP_ALIGN.CENTER
            
    # Bottom bullets
    tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.73), Inches(2.0))
    tf = tx_box.text_frame
    tf.word_wrap = True
    bullets = [
        "Zero API Dependencies: Runs completely offline, avoiding latency drift.",
        "Model Context Protocol (MCP): Unifies filesystem, shells, and debuggers."
    ]
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = "✦  " + b
        p.font.name = 'Plus Jakarta Sans'
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(10)

    # Slide 7: Why This Solution Wins
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    add_header(slide, "Why This Solution Wins", "India Runs Challenge Verification")
    
    # Left column bullets
    tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.5))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    bullets = [
        "0% Honeypot Rate: Filters out exactly 81 logical anomalies (timeline conflicts, 0-duration expert skills, future certs).",
        "Unmatched Performance: Runs in &le; 12 seconds over the 100k candidate pool.",
        "Minimal Compute Footprint: Uses ~180 MB RAM and runs CPU-only, offline.",
        "ATS Noise Reducer: Penalizes job hopping (< 18 mo) and filters consulting-only backgrounds."
    ]
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = "✦  " + b
        p.font.name = 'Plus Jakarta Sans'
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(14)
        
    # Right column box
    tx_box_r = slide.shapes.add_textbox(Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.5))
    tf_r = tx_box_r.text_frame
    tf_r.word_wrap = True
    p_r = tf_r.paragraphs[0]
    p_r.text = "STAGE 3 & 4 READY"
    p_r.font.name = 'Outfit'
    p_r.font.size = Pt(16)
    p_r.font.bold = True
    p_r.font.color.rgb = ACCENT_PRIMARY
    
    p_r2 = tf_r.add_paragraph()
    p_r2.text = "\nDesigned to meet the compute environment restrictions of Stage 3 and provide unique, non-templated candidate match reasonings for Stage 4 manual reviews."
    p_r2.font.name = 'Plus Jakarta Sans'
    p_r2.font.size = Pt(15)
    p_r2.font.color.rgb = TEXT_SECONDARY

    # Slide 8: Roadmap & Business Value
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    add_header(slide, "Roadmap & Business Value", "Future Milestones")
    
    # Left column roadmap box
    tx_box_l = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.5))
    tf_l = tx_box_l.text_frame
    tf_l.word_wrap = True
    p_l = tf_l.paragraphs[0]
    p_l.text = "ROADMAP"
    p_l.font.name = 'Outfit'
    p_l.font.size = Pt(16)
    p_l.font.bold = True
    p_l.font.color.rgb = ACCENT_SECONDARY
    
    roadmap_items = [
        "Days 1-30: Establish core docker sandboxing and compiler integration.",
        "Days 31-60: Pilot Chrome DevTools protocol UI verification with internal developers.",
        "Days 61-90: Launch v1.0 open-source VS Code MCP server extension."
    ]
    for r in roadmap_items:
        p_item = tf_l.add_paragraph()
        p_item.text = "\n" + r
        p_item.font.name = 'Plus Jakarta Sans'
        p_item.font.size = Pt(14)
        p_item.font.color.rgb = TEXT_SECONDARY
        
    # Right column bullets
    tx_box_r = slide.shapes.add_textbox(Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.5))
    tf_r = tx_box_r.text_frame
    tf_r.word_wrap = True
    bullets = [
        "70% Reduction in developer debugging and feature drafting cycles.",
        "50% Faster QA Loops due to automated visual and accessibility testing.",
        "ATS Filtering Accuracy to secure top-tier engineering talent and exclude fraud profiles."
    ]
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = "✦  " + b
        p.font.name = 'Plus Jakarta Sans'
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(20)

    # Save presentation
    prs.save("presentation.pptx")
    print("presentation.pptx generated successfully!")

if __name__ == "__main__":
    create_presentation()
